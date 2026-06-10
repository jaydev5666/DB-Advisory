from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import requests, os, json
from pymongo import MongoClient
from datetime import datetime, timedelta
import pandas as pd
from pptx import Presentation
import yfinance as yf
from dotenv import load_dotenv
import bcrypt
import jwt
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
import certifi

load_dotenv()

app = Flask(__name__)
CORS(app) # Allow cross-origin requests from frontend

API_KEY = os.getenv("API_KEY")
JWT_SECRET = os.getenv("JWT_SECRET", "super-secret-key-change-this")
ALPHA_VANTAGE_KEY = os.getenv("ALPHA_VANTAGE_KEY")
NEWS_API_KEY = os.getenv("NEWS_API_KEY")
TWELVE_DATA_KEY = os.getenv("TWELVE_DATA_KEY", "")
FINNHUB_KEY = os.getenv("FINNHUB_KEY", "")

# Rate Limiter
limiter = Limiter(
    get_remote_address,
    app=app,
    default_limits=["200 per day", "50 per hour"],
    storage_uri="memory://",
)

# MongoDB Configuration
MONGO_URI = os.getenv("MONGO_URI")
# Troubleshooting: Temporarily allow invalid certificates to check if it's a SSL or IP issue
client = MongoClient(MONGO_URI, tlsCAFile=certifi.where(), tlsAllowInvalidCertificates=True)
db = client["db_advisory"]
users_col = db["users"]
history_col = db["history"]
visitors_col = db["visitors"]
assets_col = db["assets"]
wealth_goals_col = db["wealth_goals"]
competitors_col = db["competitor_firms"]
acquisitions_col = db["acquisitions"]

def auto_seed_competitors():
    if competitors_col.count_documents({}) == 0:
        firms = [
            # Asset Managers
            {
                "firmName": "BlackRock",
                "type": "Asset Management",
                "globalRank": 1,
                "totalAUM": 11500000000000,
                "headquarters": "New York, USA",
                "primaryStrategy": "Passive Indexing & Aladdin Risk Platform",
                "techStackFocus": ["AI Risk Management", "Aladdin Risk Platform", "Institutional Portals"],
                "portfolioAllocation": { "Equities": 55, "Fixed Income": 30, "Alternatives": 10, "Cash": 5 }
            },
            {
                "firmName": "Vanguard Group",
                "type": "Asset Management",
                "globalRank": 2,
                "totalAUM": 10100000000000,
                "headquarters": "Pennsylvania, USA",
                "primaryStrategy": "Low-Cost Index Mutual Funds & ETFs",
                "techStackFocus": ["Robo-Advisory", "Automated Rebalancing", "Retail Hubs"],
                "portfolioAllocation": { "Equities": 65, "Fixed Income": 25, "Alternatives": 2, "Cash": 8 }
            },
            {
                "firmName": "Fidelity Investments",
                "type": "Asset Management",
                "globalRank": 3,
                "totalAUM": 5300000000000,
                "headquarters": "Boston, USA",
                "primaryStrategy": "Active Equity & Retail Services",
                "techStackFocus": ["Crypto Custody", "Digital Brokerage", "AI Assistants"],
                "portfolioAllocation": { "Equities": 60, "Fixed Income": 20, "Alternatives": 12, "Cash": 8 }
            },
            {
                "firmName": "State Street Global Advisors",
                "type": "Asset Management",
                "globalRank": 4,
                "totalAUM": 4300000000000,
                "headquarters": "Boston, USA",
                "primaryStrategy": "ETF Issuance & SPDR",
                "techStackFocus": ["Blockchain Settlement", "Institutional Reporting", "API Delivery"],
                "portfolioAllocation": { "Equities": 70, "Fixed Income": 20, "Alternatives": 3, "Cash": 7 }
            },
            {
                "firmName": "J.P. Morgan Asset Management",
                "type": "Asset Management",
                "globalRank": 5,
                "totalAUM": 3400000000000,
                "headquarters": "New York, USA",
                "primaryStrategy": "Multi-Asset & Active Solutions",
                "techStackFocus": ["Machine Learning Portfolio Design", "Custom Indexing", "Digital Onboarding"],
                "portfolioAllocation": { "Equities": 50, "Fixed Income": 30, "Alternatives": 15, "Cash": 5 }
            },
            {
                "firmName": "Goldman Sachs Asset Management",
                "type": "Asset Management",
                "globalRank": 6,
                "totalAUM": 2800000000000,
                "headquarters": "New York, USA",
                "primaryStrategy": "Alternative Investments & Quantitative Models",
                "techStackFocus": ["Algorithmic Trading", "AI Alpha Generators", "Institutional Analytics"],
                "portfolioAllocation": { "Equities": 40, "Fixed Income": 25, "Alternatives": 30, "Cash": 5 }
            },
            {
                "firmName": "Capital Group",
                "type": "Asset Management",
                "globalRank": 7,
                "totalAUM": 2600000000000,
                "headquarters": "Los Angeles, USA",
                "primaryStrategy": "Active Mutual Funds & Long-term Growth",
                "techStackFocus": ["Proprietary Valuation Systems", "Portfolio Modeling Tools", "Client Portal APIs"],
                "portfolioAllocation": { "Equities": 75, "Fixed Income": 15, "Alternatives": 2, "Cash": 8 }
            },
            {
                "firmName": "Amundi",
                "type": "Asset Management",
                "globalRank": 8,
                "totalAUM": 2300000000000,
                "headquarters": "Paris, France",
                "primaryStrategy": "ESG Integration & European Scale",
                "techStackFocus": ["ESG Rating Engines", "Digital Distribution Platforms", "Open Wealth APIs"],
                "portfolioAllocation": { "Equities": 45, "Fixed Income": 35, "Alternatives": 12, "Cash": 8 }
            },
            {
                "firmName": "Allianz (PIMCO)",
                "type": "Asset Management",
                "globalRank": 9,
                "totalAUM": 2200000000000,
                "headquarters": "Munich, Germany",
                "primaryStrategy": "Fixed-Income & Active Bonds",
                "techStackFocus": ["Fixed Income Risk Systems", "Macro Trade Analysis Tools", "Institutional APIs"],
                "portfolioAllocation": { "Equities": 15, "Fixed Income": 75, "Alternatives": 7, "Cash": 3 }
            },
            {
                "firmName": "Brookfield Asset Management",
                "type": "Asset Management",
                "globalRank": 10,
                "totalAUM": 1100000000000,
                "headquarters": "Toronto, Canada",
                "primaryStrategy": "Alternative Infrastructure & Tangible Assets",
                "techStackFocus": ["Renewable Grid Forecasting", "Real Estate Asset Portals", "LPs Co-investment Engine"],
                "portfolioAllocation": { "Equities": 5, "Fixed Income": 10, "Alternatives": 80, "Cash": 5 }
            },
            # Wealth Managers
            {
                "firmName": "UBS Global Wealth Management",
                "type": "Wealth Management",
                "globalRank": 1,
                "totalAUM": 6200000000000,
                "headquarters": "Zurich, Switzerland",
                "primaryStrategy": "Cross-Border Wealth & Global HNWIs",
                "techStackFocus": ["Private Wealth Digital Advisors", "Automated Wealth Structuring", "Swiss Digital Vaults"],
                "portfolioAllocation": { "Equities": 45, "Fixed Income": 30, "Alternatives": 15, "Cash": 10 }
            },
            {
                "firmName": "Morgan Stanley Wealth Management",
                "type": "Wealth Management",
                "globalRank": 2,
                "totalAUM": 5800000000000,
                "headquarters": "New York, USA",
                "primaryStrategy": "Advisor-Led Client Solutions & Digital Platform",
                "techStackFocus": ["Next Best Action AI", "Digital Wealth Platforms", "Client Advisor Bridges"],
                "portfolioAllocation": { "Equities": 50, "Fixed Income": 28, "Alternatives": 12, "Cash": 10 }
            },
            {
                "firmName": "Bank of America Merrill Lynch",
                "type": "Wealth Management",
                "globalRank": 3,
                "totalAUM": 3800000000000,
                "headquarters": "Charlotte, USA",
                "primaryStrategy": "Unified Bank & Wealth Services",
                "techStackFocus": ["Erica AI Advisor Assistant", "Mobile Wealth Management", "Retirement Benefits Engines"],
                "portfolioAllocation": { "Equities": 48, "Fixed Income": 32, "Alternatives": 10, "Cash": 10 }
            },
            {
                "firmName": "J.P. Morgan Private Bank",
                "type": "Wealth Management",
                "globalRank": 4,
                "totalAUM": 3100000000000,
                "headquarters": "New York, USA",
                "primaryStrategy": "UHNW Bespoke Planning & Lending",
                "techStackFocus": ["Bespoke Portfolio Customization", "Alternative Asset Pipelines", "UHNW Family Office Hubs"],
                "portfolioAllocation": { "Equities": 42, "Fixed Income": 25, "Alternatives": 25, "Cash": 8 }
            },
            {
                "firmName": "Charles Schwab Wealth Advisory",
                "type": "Wealth Management",
                "globalRank": 5,
                "totalAUM": 2400000000000,
                "headquarters": "Texas, USA",
                "primaryStrategy": "Self-Directed & Hybrid Advice Models",
                "techStackFocus": ["Robo-Advisory", "Zero-Commission Trading Hubs", "Financial Planning APIs"],
                "portfolioAllocation": { "Equities": 60, "Fixed Income": 25, "Alternatives": 5, "Cash": 10 }
            },
            {
                "firmName": "Goldman Sachs Private Wealth Management",
                "type": "Wealth Management",
                "globalRank": 6,
                "totalAUM": 1600000000000,
                "headquarters": "New York, USA",
                "primaryStrategy": "UHNW Advisory & Co-Investment Programs",
                "techStackFocus": ["Exclusive Alternative Asset Access", "Family Office Portals", "Risk Stress Testing Apps"],
                "portfolioAllocation": { "Equities": 35, "Fixed Income": 20, "Alternatives": 40, "Cash": 5 }
            },
            {
                "firmName": "Citi Private Bank",
                "type": "Wealth Management",
                "globalRank": 7,
                "totalAUM": 1000000000000,
                "headquarters": "New York, USA",
                "primaryStrategy": "Global HNW & Offshore Investment Access",
                "techStackFocus": ["Cross-Border Wealth Structuring", "Multi-Currency Platforms", "Offshore Tax Calculators"],
                "portfolioAllocation": { "Equities": 44, "Fixed Income": 28, "Alternatives": 18, "Cash": 10 }
            },
            {
                "firmName": "BNP Paribas Wealth Management",
                "type": "Wealth Management",
                "globalRank": 8,
                "totalAUM": 550000000000,
                "headquarters": "Paris, France",
                "primaryStrategy": "Euro-Zone Private Capital Advisory",
                "techStackFocus": ["Sustainable Investment Portals", "Private Equity Routing", "European Wealth Bridges"],
                "portfolioAllocation": { "Equities": 40, "Fixed Income": 35, "Alternatives": 15, "Cash": 10 }
            },
            {
                "firmName": "Julius Baer",
                "type": "Wealth Management",
                "globalRank": 9,
                "totalAUM": 520000000000,
                "headquarters": "Zurich, Switzerland",
                "primaryStrategy": "Pure-Play Private Banking & Swiss Wealth Management",
                "techStackFocus": ["Swiss Security Protocols", "HNW Financial Calculators", "Bespoke Portfolio Dashboards"],
                "portfolioAllocation": { "Equities": 42, "Fixed Income": 33, "Alternatives": 13, "Cash": 12 }
            },
            {
                "firmName": "HSBC Private Banking",
                "type": "Wealth Management",
                "globalRank": 10,
                "totalAUM": 480000000000,
                "headquarters": "London, UK",
                "primaryStrategy": "Asia-Focused HNW Wealth Connections",
                "techStackFocus": ["Global Wealth Transfer Portals", "International Investment Hubs", "Multi-Jurisdiction Reporting"],
                "portfolioAllocation": { "Equities": 46, "Fixed Income": 30, "Alternatives": 12, "Cash": 12 }
            }
        ]
        competitors_col.insert_many(firms)
        print("Competitor firms seeded successfully.")
    
    if acquisitions_col.count_documents({}) == 0:
        # Retrieve firm details to link correct buyers
        blackrock = competitors_col.find_one({"firmName": "BlackRock"})
        morgan_stanley = competitors_col.find_one({"firmName": "Morgan Stanley Wealth Management"})
        ubs = competitors_col.find_one({"firmName": "UBS Global Wealth Management"})
        brookfield = competitors_col.find_one({"firmName": "Brookfield Asset Management"})
        allianz = competitors_col.find_one({"firmName": "Allianz (PIMCO)"})
        jpm_am = competitors_col.find_one({"firmName": "J.P. Morgan Asset Management"})
        amundi = competitors_col.find_one({"firmName": "Amundi"})

        deals = [
            {
                "buyerFirmId": blackrock["_id"] if blackrock else None,
                "buyerFirmName": "BlackRock",
                "targetName": "Preqin",
                "targetCategory": "Fintech/Software",
                "dealValue": 3200000000,
                "status": "Completed",
                "announcementDate": datetime(2024, 7, 1),
                "strategicIntent": "Integrating private markets data and index solutions into Aladdin to enhance private assets risk stress-testing."
            },
            {
                "buyerFirmId": blackrock["_id"] if blackrock else None,
                "buyerFirmName": "BlackRock",
                "targetName": "Global Infrastructure Partners (GIP)",
                "targetCategory": "Renewable Energy",
                "dealValue": 12500000000,
                "status": "Completed",
                "announcementDate": datetime(2024, 1, 12),
                "strategicIntent": "Scaling multi-asset infrastructure investing in utility solar, digital networks, and logistics networks globally."
            },
            {
                "buyerFirmId": morgan_stanley["_id"] if morgan_stanley else None,
                "buyerFirmName": "Morgan Stanley Wealth Management",
                "targetName": "E*TRADE Financial",
                "targetCategory": "Competitor Firm",
                "dealValue": 13000000000,
                "status": "Completed",
                "announcementDate": datetime(2020, 2, 20),
                "strategicIntent": "Acquiring a digital direct-to-consumer brokerage channel to feed downstream advisory and stock plan systems."
            },
            {
                "buyerFirmId": morgan_stanley["_id"] if morgan_stanley else None,
                "buyerFirmName": "Morgan Stanley Wealth Management",
                "targetName": "Eaton Vance",
                "targetCategory": "Competitor Firm",
                "dealValue": 7000000000,
                "status": "Completed",
                "announcementDate": datetime(2020, 10, 8),
                "strategicIntent": "Enhancing custom indexation capabilities (Parametric) and active fixed-income assets under management."
            },
            {
                "buyerFirmId": ubs["_id"] if ubs else None,
                "buyerFirmName": "UBS Global Wealth Management",
                "targetName": "Credit Suisse",
                "targetCategory": "Competitor Firm",
                "dealValue": 3200000000,
                "status": "Completed",
                "announcementDate": datetime(2023, 3, 19),
                "strategicIntent": "Consolidating Swiss private banking leadership and creating a global private wealth giant managing over $5 Trillion."
            },
            {
                "buyerFirmId": brookfield["_id"] if brookfield else None,
                "buyerFirmName": "Brookfield Asset Management",
                "targetName": "Neoen (French Renewable Giant)",
                "targetCategory": "Renewable Energy",
                "dealValue": 6500000000,
                "status": "Exclusive Talks",
                "announcementDate": datetime(2024, 5, 30),
                "strategicIntent": "Acquiring 8 GW of wind, solar, and battery storage assets to accelerate decarbonization strategies."
            },
            {
                "buyerFirmId": brookfield["_id"] if brookfield else None,
                "buyerFirmName": "Brookfield Asset Management",
                "targetName": "American Landmark Apartments",
                "targetCategory": "Real Estate",
                "dealValue": 1500000000,
                "status": "Completed",
                "announcementDate": datetime(2023, 8, 15),
                "strategicIntent": "Expanding multi-family residential portfolio across high-growth Sun Belt metros in the United States."
            },
            {
                "buyerFirmId": allianz["_id"] if allianz else None,
                "buyerFirmName": "Allianz (PIMCO)",
                "targetName": "UOB Asset Management (Singapore)",
                "targetCategory": "Competitor Firm",
                "dealValue": 150000000,
                "status": "Exclusive Talks",
                "announcementDate": datetime(2024, 4, 10),
                "strategicIntent": "Gaining local Southeast Asian asset management license and distribution channels through UOB's branch networks."
            },
            {
                "buyerFirmId": jpm_am["_id"] if jpm_am else None,
                "buyerFirmName": "J.P. Morgan Asset Management",
                "targetName": "Nutmeg (UK Robo-Advisor)",
                "targetCategory": "Fintech/Software",
                "dealValue": 1000000000,
                "status": "Completed",
                "announcementDate": datetime(2021, 6, 17),
                "strategicIntent": "Launching digital wealth offerings in the UK and Europe as part of Chase international retail banking rollout."
            },
            {
                "buyerFirmId": amundi["_id"] if amundi else None,
                "buyerFirmName": "Amundi",
                "targetName": "Lyxor Asset Management",
                "targetCategory": "Competitor Firm",
                "dealValue": 975000000,
                "status": "Completed",
                "announcementDate": datetime(2021, 4, 7),
                "strategicIntent": "Creating the second largest ETF provider in Europe, boosting passive assets scale to compete with US indexing giants."
            }
        ]
        acquisitions_col.insert_many(deals)
        print("Acquisitions seeded successfully.")

def init_db():
    # Add default admin user if not exists and ensure admin role from .env
    admin_email = os.getenv("ADMIN_EMAIL", "admin")
    admin_password = os.getenv("ADMIN_PASSWORD", "1234")
    
    # Try to find a user with role "admin"
    admin_user = users_col.find_one({"role": "admin"})
    if not admin_user:
        # Fallback to finding user by username "admin" to migrate
        admin_user = users_col.find_one({"username": "admin"})
        
    if not admin_user:
        hashed = bcrypt.hashpw(admin_password.encode('utf-8'), bcrypt.gensalt())
        users_col.insert_one({
            "username": admin_email,
            "password": hashed,
            "role": "admin",
            "created_at": datetime.utcnow()
        })
    else:
        updates = {}
        if admin_user.get("username") != admin_email:
            updates["username"] = admin_email
        if admin_user.get("role") != "admin":
            updates["role"] = "admin"
            
        try:
            if not bcrypt.checkpw(admin_password.encode('utf-8'), admin_user['password']):
                updates["password"] = bcrypt.hashpw(admin_password.encode('utf-8'), bcrypt.gensalt())
        except Exception:
            updates["password"] = bcrypt.hashpw(admin_password.encode('utf-8'), bcrypt.gensalt())
            
        if updates:
            users_col.update_one({"_id": admin_user["_id"]}, {"$set": updates})
            
    print("MongoDB connection established and initialized.")
    auto_seed_competitors()

init_db()

def get_fallback_data(company, deal_type):
    # Try to clean company name/ticker
    search_term = company.strip().upper()
    
    try:
        stock = yf.Ticker(search_term)
        info = stock.info
        
        # If no info, try to search for the ticker first
        if not info or 'symbol' not in info or 'sector' not in info:
            print(f"No direct info for {search_term}, trying search...")
            # Newer yfinance versions use yf.Search().tickers but older ones might differ
            try:
                search = yf.Search(search_term).tickers
                if search and len(search) > 0:
                    search_term = search[0]
                    stock = yf.Ticker(search_term)
                    info = stock.info
            except Exception:
                # Fallback search method
                ticker = yf.Ticker(search_term)
                if not ticker.info:
                    print("Ticker search failed.")
                else:
                    info = ticker.info

        if search_term.lower() in ["g", "genpact"]:
            return {
                "overview": "Status: Publicly traded on the NYSE (Ticker: G) since its IPO in August 2007. Current Price: $34.14 (as of May 4, 2026). Market Cap: Approximately $5.76B - $5.83B.",
                "industry": "Sector: Professional Services / IT & Business Process Management (BPM). Key Trend: Transitioning to Agentic AI and 'AI Gigafactories.' 92% of executives in Genpact's network expect agentic AI to fundamentally change operations by 2027.",
                "peers": [], # Let AI decide or leave empty if fallback
                "valuation": "P/E Ratio (TTM): 10.97x, which is currently lower than its 5-year average, reflecting sector-wide 'AI deflation' concerns. Dividend Yield: 2.18% (Annual dividend of $0.75 per share). Revenue Growth: 2026 outlook projects at least 7% reported growth.",
                "rationale": "AI Leadership: Recognized as a 'Leader' in the ISG Provider Lens 2026 for ServiceNow ecosystem partners. Shareholder Returns: Aggressive dividend increases ($0.68 to $0.75) and high institutional ownership (~96%). Strategic Positioning: Strong pivot from traditional BPO to AI-led process intelligence.",
                "risks": "AI Deflation: General market fear that AI efficiency gains might reduce total contract values for traditional outsourcing. Technical Weakness: Shares are currently trading near 52-week lows ($33.14) and below 50/200-day moving averages. Earnings Volatility: Q1 2026 earnings are scheduled for May 7, 2026.",
                "banks": [
                    {"name": "Goldman Sachs", "ecm_score": 5, "sector_score": 5, "region_score": 5},
                    {"name": "Morgan Stanley", "ecm_score": 5, "sector_score": 4, "region_score": 5}
                ],
                "figures": "Market Cap: ~$5.8B, Current Price: $34.14, P/E: 10.97x, Dividend Yield: 2.18%."
            }
        
        if info and "sector" in info:
            name = info.get("shortName", search_term)
            sector = info.get("sector", "N/A")
            industry = info.get("industry", "N/A")
            summary = info.get("longBusinessSummary", "")
            market_cap = info.get("marketCap", "N/A")
            
            # Detect currency symbol dynamically
            u_term = search_term.upper()
            if any(x in u_term for x in ['.NS', '.BO']):
                sym = '₹'
            elif '.L' in u_term:
                sym = '£'
            elif any(x in u_term for x in ['.DE', '.PA', '.AS', '.MI']):
                sym = '€'
            elif '.T' in u_term:
                sym = '¥'
            else:
                sym = '$'
                
            if isinstance(market_cap, (int, float)): market_cap = f"{sym}{market_cap / 1e9:.2f}B"
            price = info.get("currentPrice", "N/A")
            pe = info.get("trailingPE", "N/A")
            
            return {
                "overview": f"{name} ({search_term}) is a publicly traded company. Current Price: {sym}{price}. Market Cap: {market_cap}.",
                "industry": f"Sector: {sector} / Industry: {industry}. {summary}",
                "peers": [],
                "valuation": f"P/E Ratio: {pe}. Valuation driven by {industry} market trends.",
                "rationale": f"Based on its position in the {sector} sector and a market cap of {market_cap}, {name} represents a standard market opportunity.",
                "risks": f"General market volatility, sector-specific macroeconomic shifts in {industry}, and competitive pressures.",
                "banks": [
                    {"name": "Goldman Sachs", "ecm_score": 5, "sector_score": 5, "region_score": 5},
                    {"name": "Morgan Stanley", "ecm_score": 4, "sector_score": 5, "region_score": 4},
                    {"name": "JP Morgan", "ecm_score": 5, "sector_score": 4, "region_score": 5}
                ],
                "figures": f"Price: {sym}{price}, Market Cap: {market_cap}, P/E: {pe}"
            }
    except Exception as e:
        print(f"YFinance fallback error: {e}")
        
    return {
        "overview": f"Detailed {deal_type} intelligence for {company} based on current market positioning.",
        "industry": "Institutional Research",
        "peers": [],
        "valuation": "Analysis based on historical sector multiples.",
        "rationale": "Strategic growth potential and market leadership.",
        "risks": "Standard macroeconomic and operational risks apply.",
        "banks": [
            {"name": "Tier 1 Bank", "ecm_score": 5, "sector_score": 5, "region_score": 5},
            {"name": "Tier 2 Bank", "ecm_score": 4, "sector_score": 4, "region_score": 4}
        ],
        "figures": "Real-time figures currently unavailable for this specific search."
    }

def get_alpha_vantage(ticker):
    if not ALPHA_VANTAGE_KEY: return None
    try:
        # Global Quote for price
        quote_url = f"https://www.alphavantage.co/query?function=GLOBAL_QUOTE&symbol={ticker}&apikey={ALPHA_VANTAGE_KEY}"
        quote_data = requests.get(quote_url).json().get("Global Quote", {})
        
        # Overview for fundamentals
        ov_url = f"https://www.alphavantage.co/query?function=OVERVIEW&symbol={ticker}&apikey={ALPHA_VANTAGE_KEY}"
        ov_data = requests.get(ov_url).json()
        
        return {
            "price": quote_data.get("05. price", "N/A"),
            "change": quote_data.get("09. change", "N/A"),
            "market_cap": ov_data.get("MarketCapitalization", "N/A"),
            "pe_ratio": ov_data.get("PERatio", "N/A"),
            "ev_ebitda": ov_data.get("EVToEBITDA", "N/A"),
            "revenue_growth": ov_data.get("QuarterlyRevenueGrowthYOY", "N/A"),
            "description": ov_data.get("Description", "")
        }
    except Exception as e:
        print(f"Alpha Vantage error: {e}")
        return None

def get_news_sentiment(company):
    if not NEWS_API_KEY: return {"headlines": [], "sentiment": "Neutral"}
    try:
        import re
        query_term = company
        
        # 1. Resolve company name if it's a ticker or has exchange suffixes to ensure high relevancy
        if company and company != "finance AND market":
            # Strip exchange suffix (e.g. TCS.NS -> TCS)
            if "." in query_term:
                parts = query_term.split(".")
                if len(parts[-1]) <= 3:
                    query_term = ".".join(parts[:-1])
            
            # Fetch company name from yfinance if query looks like a ticker (short and alphanumeric, or has dots)
            clean_term = query_term.strip()
            if ("." in company) or (len(clean_term) <= 12 and clean_term.replace(".", "").replace("-", "").isalnum()):
                try:
                    ticker_obj = yf.Ticker(company)
                    info = ticker_obj.info
                    resolved_name = info.get("longName") or info.get("shortName")
                    if resolved_name:
                        # Clean legal suffixes case-insensitively using regex word boundaries to avoid corruption
                        # e.g., Inc., Ltd., Corporation, LLC, PLC, etc.
                        cleaned = resolved_name
                        for _ in range(3): # pass up to 3 times to clean nested suffixes (e.g., Co., Ltd.)
                            prev = cleaned
                            cleaned = re.sub(
                                r'\b(inc|ltd|corp|corporation|limited|co|company|llc|plc|pty)\b\.?$', 
                                '', 
                                cleaned, 
                                flags=re.IGNORECASE
                            ).strip()
                            cleaned = re.sub(r'[\s,]+$', '', cleaned).strip()
                            if cleaned == prev:
                                break
                        query_term = cleaned
                except Exception as yf_err:
                    print(f"yfinance helper error in news query: {yf_err}")
        
        # 2. Build strict query to filter out irrelevant general news on dashboard
        search_query = query_term.strip()
        if search_query and search_query != "finance AND market":
            search_query = f'"{search_query}" AND (finance OR stock OR earnings OR business OR deal OR acquisition OR merger)'
        else:
            search_query = "finance AND market"

        # 3. Query NewsAPI safely using request params to handle encoding
        url = "https://newsapi.org/v2/everything"
        params = {
            "q": search_query,
            "sortBy": "publishedAt",
            "pageSize": 5,
            "apiKey": NEWS_API_KEY
        }
        articles = requests.get(url, params=params).json().get("articles", [])
        
        headlines = [{"title": a["title"], "source": a["source"]["name"], "url": a["url"]} for a in articles]
        
        bullish_words = ["growth", "surge", "upgraded", "profit", "beats", "expansion", "buy"]
        bearish_words = ["drop", "slump", "down", "loss", "misses", "risks", "sell"]
        
        score = 0
        for h in headlines:
            text = h["title"].lower()
            for w in bullish_words: 
                if w in text: score += 1
            for w in bearish_words: 
                if w in text: score -= 1
        
        sentiment = "Bullish" if score > 1 else ("Bearish" if score < -1 else "Neutral")
        return {"headlines": headlines, "sentiment": sentiment, "score": score}
    except Exception as e:
        print(f"NewsAPI error: {e}")
        return {"headlines": [], "sentiment": "Neutral"}

def merge_market_context(company):
    av_data = get_alpha_vantage(company)
    news_data = get_news_sentiment(company)
    
    context = f"COMPANY: {company}\n"
    if av_data:
        context += f"FINANCIALS: Price: {av_data['price']}, P/E: {av_data['pe_ratio']}, EV/EBITDA: {av_data['ev_ebitda']}, Growth: {av_data['revenue_growth']}\n"
        if av_data['description']: context += f"SUMMARY: {av_data['description'][:300]}...\n"
    
    if news_data['headlines']:
        h_text = " | ".join([h['title'] for h in news_data['headlines']])
        context += f"RECENT NEWS: {h_text}\n"
        context += f"MARKET SENTIMENT: {news_data['sentiment']}\n"
        
    return context, av_data, news_data

def generate_analysis(company, deal_type, market_context=""):
    if not API_KEY or API_KEY == "sk-proj-your-real-api-key-goes-here...":
        return get_fallback_data(company, deal_type)
        
    prompt = f"""
Act as a Senior Investment Banking Analyst at a Tier 1 global investment bank.

Today's Date: {datetime.utcnow().strftime('%Y-%m-%d')}
Company: {company}
Target Deal Type: {deal_type}

LIVE MARKET CONTEXT (MANDATORY TO USE):
{market_context}

STRICT ANALYSIS INSTRUCTIONS:

1. DEAL CONTEXT ENFORCEMENT:
   - If Deal Type is "IPO": Focus on capital structure, listing prospects, equity story, and potential valuation ranges. If the company is already public, treat this as a "Secondary/Follow-on Offering" and focus on how the new capital will be used.
   - If Deal Type is "M&A": Focus on strategic fit, potential synergies (cost and revenue), market consolidation, and valuation premium.
   - If Deal Type is "LBO": Focus on cash flow stability, debt capacity, operational improvement opportunities, and exit strategies (e.g., IPO or Trade Sale).
   - If Deal Type is "Restructuring": Focus on debt re-profiling, asset sales, liquidity management, and turnaround strategy.

2. Determine whether {company} is already publicly traded. If public, explicitly mention the current stock ticker and exchange in the overview.

3. Use the most recent available market intelligence, earnings trends, valuation data, sector performance, and capital markets activity from 2024/2025 knowledge.

4. The response MUST sound like a real institutional equity research / ECM banking note:
   - Highly analytical
   - Numerically driven
   - Concise but insight-heavy
   - No generic wording
   - No placeholder language
   - Include real market commentary and sector sentiment

5. VERY IMPORTANT:
   - Blend financial figures naturally INTO the text.
   - Do NOT output robotic isolated numbers.
   - Example style:
     "The company generated approximately $4.5B revenue with EBITDA margins near 18%, while trading around 16x forward EBITDA versus peers at 18-22x."
   - Every section should feel like authentic sell-side banking commentary.

6. Word Limit Guidance:
   - "overview" and "industry" MUST be comprehensive (150-250 words each).
   - "valuation" should be a detailed 200-word analysis of multiples and market positioning.
   - Depth and quality are CRITICAL. Do not truncate or provide short summaries.

Return ONLY a valid JSON object with EXACTLY these keys:

{{
  "overview": "Professional executive summary with ticker (if public), market positioning, strategic context, and transaction framing.",
  
  "industry": "Detailed industry analysis including sector growth, AI/digital trends, recent M&A activity, macroeconomic headwinds, margin trends, and investor sentiment.",
  
  "peers": ["TICKER1", "TICKER2", "TICKER3"],
  
  "valuation": "Detailed valuation commentary using current EV/EBITDA, P/E, revenue multiples, peer premium/discount discussion, and public market positioning.",
  
  "rationale": [
    "Investment rationale point 1",
    "Investment rationale point 2",
    "Investment rationale point 3",
    "Investment rationale point 4"
  ],
  
  "risks": [
    "Commercial/financial risk 1",
    "Commercial/financial risk 2",
    "Commercial/financial risk 3",
    "Commercial/financial risk 4"
  ],
  
  "banks": [
    {{
      "name": "Bank Name",
      "ecm_score": 1-5,
      "sector_score": 1-5,
      "region_score": 1-5
    }}
  ],
  
  "figures": "Integrated professional financial summary including Market Cap, Revenue, EBITDA, margins, growth profile, leverage commentary, valuation range, and stock performance where relevant."
}}

FINAL OUTPUT RULES:
- Output ONLY valid JSON.
- No markdown.
- No explanations.
- No introductory text.
- No trailing comments.
"""

    # Determine which endpoint to use
    is_openrouter = API_KEY.startswith("sk-or")
    is_gemini = API_KEY.startswith("AIzaSy")
    
    try:
        if is_gemini:
            # Use stable v1 API
            url = f"https://generativelanguage.googleapis.com/v1/models/gemini-1.5-flash:generateContent?key={API_KEY}"
            payload = {
                "contents": [{
                    "parts": [{"text": prompt + "\nRespond ONLY with valid JSON. Do NOT include markdown code blocks or backticks."}]
                }]
            }
            response = requests.post(url, json=payload)
            res_json = response.json()
            
            if 'error' in res_json:
                print(f"Gemini API Error Detail:", json.dumps(res_json, indent=2))
                return get_fallback_data(company, deal_type)
            
            if 'candidates' not in res_json or not res_json['candidates']:
                print(f"Gemini Response Error: No candidates found. Response: {json.dumps(res_json, indent=2)}")
                return get_fallback_data(company, deal_type)
                
            content = res_json['candidates'][0]['content']['parts'][0]['text']
            # Clean up potential markdown backticks from Gemini
            content = content.replace("```json", "").replace("```", "").strip()
            return json.loads(content)

        url = "https://openrouter.ai/api/v1/chat/completions" if is_openrouter else "https://api.openai.com/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json"
        }
        if is_openrouter:
            headers["HTTP-Referer"] = "http://localhost:5174"
            headers["X-Title"] = "DB Advisory Platform"

        response = requests.post(
            url,
            headers=headers,
            json={
                "model": "openai/gpt-4o-mini" if is_openrouter else "gpt-4o-mini",
                "messages": [{"role": "user", "content": prompt}],
                "response_format": {"type": "json_object"}
            }
        )
        res_json = response.json()
        if 'error' in res_json:
            print(f"API Error Detail:", json.dumps(res_json, indent=2))
            return get_fallback_data(company, deal_type)
            
        content = res_json['choices'][0]['message']['content']
        return json.loads(content)
    except Exception as e:
        print("API CALL CRITICAL ERROR:", e)
        import traceback
        traceback.print_exc()
        return get_fallback_data(company, deal_type)

def get_historical_data(company_query, period="1y"):
    try:
        period_map = {
            "1D": {"period": "1d", "interval": "5m", "format": "%I:%M %p"},
            "1W": {"period": "5d", "interval": "1h", "format": "%a %I %p"},
            "3M": {"period": "3mo", "interval": "1d", "format": "%b %d"},
            "6M": {"period": "6mo", "interval": "1d", "format": "%b %d"},
            "1Y": {"period": "1y", "interval": "1d", "format": "%b %Y"},
            "5Y": {"period": "5y", "interval": "1wk", "format": "%Y"},
        }
        p_info = period_map.get(period.upper(), period_map["1Y"])

        ticker = company_query
        try:
            search_res = requests.get(f"https://query2.finance.yahoo.com/v1/finance/search?q={company_query}", headers={'User-Agent': 'Mozilla/5.0'}, timeout=5).json()
            if search_res.get('quotes'):
                ticker = search_res['quotes'][0]['symbol']
        except:
            pass

        # Detect native currency from ticker suffix (fast, no extra API call)
        t_upper = ticker.upper()
        if any(x in t_upper for x in ['.NS', '.BO']):
            stock_currency = 'INR'
        elif '.L' in t_upper:
            stock_currency = 'GBP'
        elif any(x in t_upper for x in ['.DE', '.PA', '.AS', '.MI']):
            stock_currency = 'EUR'
        elif '.T' in t_upper:
            stock_currency = 'JPY'
        else:
            stock_currency = 'USD'

        stock = yf.Ticker(ticker)
        hist = stock.history(period=p_info["period"], interval=p_info["interval"])
        if hist.empty:
            return {"data": [], "currency": stock_currency}

        max_points = 50 if period in ["1D", "1W"] else 30
        if len(hist) > max_points:
            step = len(hist) // max_points
            hist = hist.iloc[::step]

        chart_data = []
        for date, row in hist.iterrows():
            try:
                date_str = date.strftime(p_info["format"])
            except:
                date_str = str(date)
            chart_data.append({"date": date_str, "price": round(row['Close'], 2)})

        return {"data": chart_data, "currency": stock_currency}
    except Exception as e:
        print(f"History fetch error for {company_query}: {e}")
        return {"data": [], "currency": "USD"}


def get_financials(company_query):
    try:
        ticker = company_query
        search_res = requests.get(f"https://query2.finance.yahoo.com/v1/finance/search?q={company_query}", headers={'User-Agent': 'Mozilla/5.0'}).json()
        if search_res.get('quotes'):
            ticker = search_res['quotes'][0]['symbol']

        stock = yf.Ticker(ticker)
        info = stock.info
        return {
            "name": info.get("shortName", ticker),
            "price": info.get("currentPrice", 0),
            "pe_ratio": info.get("trailingPE", 0),
            "ev_ebitda": info.get("enterpriseToEbitda", 0)
        }
    except Exception as e:
        return {"name": company_query, "price": 0, "pe_ratio": 0, "ev_ebitda": 0}

def generate_comps(peers):
    if not peers or not isinstance(peers, list):
        return []
    
    comps = [["Company", "Price", "P/E", "EV/EBITDA"]]
    for p in peers[:3]: 
        data = get_financials(p)
        
        # Detect currency symbol dynamically
        p_upper = p.upper()
        if any(x in p_upper for x in ['.NS', '.BO']):
            sym = '₹'
        elif '.L' in p_upper:
            sym = '£'
        elif any(x in p_upper for x in ['.DE', '.PA', '.AS', '.MI']):
            sym = '€'
        elif '.T' in p_upper:
            sym = '¥'
        else:
            sym = '$'
            
        price = f"{sym}{data['price']}" if isinstance(data['price'], (int, float)) else data['price']
        pe = f"{data['pe_ratio']:.1f}x" if isinstance(data['pe_ratio'], (int, float)) else data['pe_ratio']
        ev = f"{data['ev_ebitda']:.1f}x" if isinstance(data['ev_ebitda'], (int, float)) else data['ev_ebitda']
        comps.append([p, price, pe, ev])
    return comps

def generate_bank_matrix(company, deal_type, ai_banks=None):
    # Use AI-provided banks if available, otherwise use defaults
    if ai_banks and isinstance(ai_banks, list) and len(ai_banks) > 0:
        banks = []
        for b in ai_banks:
            banks.append({
                "name": b.get("name", "Unknown Bank"),
                "ecm": b.get("ecm_score", 3),
                "sector": b.get("sector_score", 3),
                "region": b.get("region_score", 3)
            })
    else:
        banks = [
            {"name": "Goldman Sachs", "ecm": 5, "sector": 5, "region": 5},
            {"name": "Morgan Stanley", "ecm": 5, "sector": 4, "region": 5},
            {"name": "JP Morgan", "ecm": 4, "sector": 4, "region": 5},
            {"name": "Kotak Mahindra Capital", "ecm": 4, "sector": 4, "region": 3},
            {"name": "ICICI Securities", "ecm": 3, "sector": 3, "region": 4}
        ]
        
    results = []
    for b in banks:
        score = (b["ecm"] * 0.4) + (b["sector"] * 0.3) + (b["region"] * 0.3)
        if score >= 4.5: prob = "High"
        elif score >= 3.5: prob = "Medium"
        else: prob = "Low"
        results.append({
            "bank": b["name"],
            "score": round(score,2),
            "probability": prob
        })
    return results

def get_logged_in_user():
    auth_header = request.headers.get("Authorization")
    if not auth_header:
        return None
    try:
        if auth_header.startswith("Bearer "):
            token = auth_header.split(" ")[1]
        else:
            token = auth_header
        data = jwt.decode(token, JWT_SECRET, algorithms=["HS256"])
        return users_col.find_one({"username": data.get("username")})
    except Exception as e:
        print(f"JWT decode error: {e}")
        return None

@app.route('/track-visit', methods=['POST'])
def track_visit():
    ip = request.headers.get('X-Forwarded-For', request.remote_addr)
    if ip and ',' in ip:
        ip = ip.split(',')[0].strip()
    
    today = datetime.utcnow().strftime("%Y-%m-%d")
    
    existing = visitors_col.find_one({"ip": ip, "date": today})
    if not existing:
        visitors_col.insert_one({
            "ip": ip,
            "date": today,
            "timestamp": datetime.utcnow()
        })
        return jsonify({"status": "success", "new_visit": True})
    return jsonify({"status": "success", "new_visit": False})

@app.route('/login', methods=['POST'])
def login():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    
    user = users_col.find_one({"username": username})
    
    if user and bcrypt.checkpw(password.encode('utf-8'), user['password']):
        token = jwt.encode({
            'username': username,
            'exp': datetime.utcnow() + timedelta(hours=24)
        }, JWT_SECRET, algorithm="HS256")
        
        return jsonify({
            "status": "success",
            "token": token,
            "user": {"username": username, "name": user.get('name', username), "role": user.get('role', 'user')}
        })
    return jsonify({"status": "fail", "message": "Invalid credentials"})

@app.route('/signup', methods=['POST'])
def signup():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    
    if not username or not password:
        return jsonify({"status": "fail", "message": "Missing credentials"})
    
    if users_col.find_one({"username": username}):
        return jsonify({"status": "fail", "message": "Username already exists"})
    
    hashed = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt())
    users_col.insert_one({"username": username, "password": hashed, "created_at": datetime.utcnow()})
    return jsonify({"status": "success", "message": "Account created successfully"})

@app.route('/google-login', methods=['POST'])
def google_login():
    data = request.json
    username = data.get('username') # This will be the email
    name = data.get('name')
    
    # Check if user exists
    user = users_col.find_one({"username": username})
    
    if not user:
        # Create new user for Google login (no password needed)
        user_data = {
            "username": username,
            "name": name,
            "auth_type": "google",
            "created_at": datetime.utcnow()
        }
        users_col.insert_one(user_data)
        user = user_data
    
    token = jwt.encode({
        'username': username,
        'exp': datetime.utcnow() + timedelta(hours=24)
    }, JWT_SECRET, algorithm="HS256")
    
    return jsonify({
        "status": "success", 
        "token": token,
        "user": {"username": username, "name": name or username, "role": user.get('role', 'user')}
    })

@app.route('/news', methods=['GET'])
def get_news():
    company = request.args.get('company')
    if not company:
        company = "finance AND market" # Default to general market news if no company provided
    news = get_news_sentiment(company)
    return jsonify(news)

@app.route('/live-quote', methods=['GET'])
def get_live_quote():
    ticker = request.args.get('ticker')
    if not ticker: return jsonify({"error": "Missing ticker parameter"}), 400
    quote = get_alpha_vantage(ticker)
    if not quote:
        # Fallback to yfinance for price only
        try:
            stock = yf.Ticker(ticker)
            price = stock.info.get('currentPrice', 'N/A')
            return jsonify({"price": price, "source": "yfinance"})
        except:
            return jsonify({"error": "Failed to fetch quote"}), 500
    quote["source"] = "Alpha Vantage"
    return jsonify(quote)

@app.route('/chart-data', methods=['GET'])
def fetch_chart_data():
    company = request.args.get('company')
    period = request.args.get('period', '1Y')
    if not company: return jsonify({"error": "Missing company parameter"}), 400
    result = get_historical_data(company, period)
    # Always return {data, currency} format
    if isinstance(result, list):
        result = {"data": result, "currency": "USD"}
    return jsonify(result)

@app.route('/analyze', methods=['POST'])
@limiter.limit("5 per minute")
def analyze():
    logged_user = get_logged_in_user()
    username = logged_user["username"] if logged_user else "anonymous"

    data = request.json
    company = data.get('company')
    deal_type = data.get('deal_type')
    
    # Fetch live figures and news for context
    market_context, av_data, news_data = merge_market_context(company)
    
    # If Alpha Vantage fails, try yfinance as fallback for context
    if not market_context or "FINANCIALS" not in market_context:
        try:
            stock = yf.Ticker(company)
            info = stock.info
            if info:
                market_context += f"\nFALLBACK FINANCIALS: Market Cap: {info.get('marketCap', 'N/A')}, Revenue: {info.get('totalRevenue', 'N/A')}, EBITDA: {info.get('ebitda', 'N/A')}, Trailing P/E: {info.get('trailingPE', 'N/A')}"
        except:
            pass

    analysis = generate_analysis(company, deal_type, market_context)
    peers = analysis.get("peers", []) if isinstance(analysis, dict) else []
    ai_banks = analysis.get("banks", []) if isinstance(analysis, dict) else []
    
    comps = generate_comps(peers)
    banks = generate_bank_matrix(company, deal_type, ai_banks)
    
    # Use new return format {data, currency}
    result = get_historical_data(company)
    history_data = result["data"] if isinstance(result, dict) else result
    benchmarking_data = []

    # Add target company first
    target_stats = get_financials(company)
    benchmarking_data.append(target_stats)
    # Add peers
    for p in peers:
        p_stats = get_financials(p)
        benchmarking_data.append(p_stats)

    result_data = {
        "analysis": analysis,
        "comps": comps,
        "banks": banks,
        "company": company,
        "deal_type": deal_type,
        "history_data": history_data,
        "currency": result.get("currency", "USD") if isinstance(result, dict) else "USD",
        "benchmarking_data": benchmarking_data,
        "news_data": news_data,
        "live_market_data": av_data,
        "username": username,
        "timestamp": datetime.utcnow()
    }
    
    # Save to MongoDB
    history_col.insert_one(result_data)
    
    # Remove _id before returning (not JSON serializable)
    if "_id" in result_data:
        del result_data["_id"]
        
    return jsonify(result_data)

@app.route('/history', methods=['GET'])
def get_history():
    rows = list(history_col.find().sort("timestamp", -1).limit(50))
    for row in rows:
        if "_id" in row:
            del row["_id"]
        # Convert datetime to string for JSON serialization
        if "timestamp" in row:
            row["timestamp"] = row["timestamp"].isoformat()
            
    return jsonify(rows)

@app.route('/admin/stats', methods=['GET'])
def admin_stats():
    logged_user = get_logged_in_user()
    if not logged_user or logged_user.get("role") != "admin":
        return jsonify({"error": "Unauthorized"}), 403
    
    try:
        total_visits = visitors_col.count_documents({})
        unique_ips_pipeline = [{"$group": {"_id": "$ip"}}, {"$count": "count"}]
        unique_ips_res = list(visitors_col.aggregate(unique_ips_pipeline))
        total_unique_visitors = unique_ips_res[0]["count"] if unique_ips_res else 0
        
        total_users = users_col.count_documents({})
        total_deals = history_col.count_documents({})
        
        recent_deals_cursor = history_col.find().sort("timestamp", -1).limit(10)
        recent_deals = []
        for deal in recent_deals_cursor:
            recent_deals.append({
                "company": deal.get("company", "N/A"),
                "deal_type": deal.get("deal_type", "N/A"),
                "username": deal.get("username", "anonymous"),
                "timestamp": deal.get("timestamp").isoformat() if deal.get("timestamp") else None
            })
            
        users_cursor = users_col.find()
        users_list = []
        for user in users_cursor:
            u_name = user.get("username")
            search_count = history_col.count_documents({"username": u_name})
            users_list.append({
                "username": u_name,
                "name": user.get("name", u_name),
                "role": user.get("role", "user"),
                "created_at": user.get("created_at").isoformat() if user.get("created_at") else None,
                "searches_count": search_count
            })
            
        return jsonify({
            "status": "success",
            "metrics": {
                "total_unique_visitors": total_unique_visitors,
                "total_visits": total_visits,
                "total_users": total_users,
                "total_deals": total_deals
            },
            "recent_deals": recent_deals,
            "users": users_list
        })
    except Exception as e:
        print(f"Error fetching admin stats: {e}")
        return jsonify({"error": "Failed to load admin stats", "details": str(e)}), 500

@app.route('/download_ppt', methods=['POST'])
def download_ppt():
    data = request.json
    prs = Presentation()
    
    def add_slide(title, content):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if hasattr(slide.placeholders[1], 'text'):
            slide.placeholders[1].text = content

    analysis = data.get('analysis', {})
    
    add_slide("Deal Overview", analysis.get('overview', 'N/A'))
    add_slide("Industry Landscape", analysis.get('industry', 'N/A'))
    
    comps_str = "Comparable Companies:\n"
    for r in data.get('comps', []):
        comps_str += " | ".join(str(x) for x in r) + "\n"
    add_slide("Valuation Comps", comps_str)
    
    add_slide("Investment Rationale", analysis.get('rationale', 'N/A'))
    add_slide("Key Risks", analysis.get('risks', 'N/A'))
    
    file_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "pitchbook.pptx"))
    prs.save(file_path)
    
    return send_file(file_path, as_attachment=True, download_name="pitchbook.pptx")


# ── Twelve Data / Finnhub Market Routes ──────────────────────

@app.route('/market/quote', methods=['GET'])
def market_quote():
    """Real-time quote for any NSE/NYSE stock via Twelve Data, falls back to yfinance."""
    symbol   = request.args.get('symbol', 'RELIANCE')
    exchange = request.args.get('exchange', 'NSE')
    try:
        if TWELVE_DATA_KEY:
            url = f"https://api.twelvedata.com/quote?symbol={symbol}&exchange={exchange}&apikey={TWELVE_DATA_KEY}"
            resp = requests.get(url, timeout=6).json()
            if resp.get('status') != 'error' and 'close' in resp:
                return jsonify({
                    "symbol":  resp.get('symbol', symbol),
                    "price":   resp.get('close', 'N/A'),
                    "open":    resp.get('open', 'N/A'),
                    "high":    resp.get('high', 'N/A'),
                    "low":     resp.get('low', 'N/A'),
                    "change":  resp.get('change', 'N/A'),
                    "pct_change": resp.get('percent_change', 'N/A'),
                    "volume":  resp.get('volume', 'N/A'),
                    "source":  "Twelve Data"
                })
        # Fallback: yfinance
        yf_sym = f"{symbol}.NS" if exchange == 'NSE' else symbol
        stock = yf.Ticker(yf_sym)
        info = stock.info
        return jsonify({
            "symbol":  symbol,
            "price":   info.get('currentPrice', info.get('regularMarketPrice', 'N/A')),
            "open":    info.get('regularMarketOpen', 'N/A'),
            "high":    info.get('dayHigh', 'N/A'),
            "low":     info.get('dayLow', 'N/A'),
            "change":  'N/A',
            "pct_change": info.get('regularMarketChangePercent', 'N/A'),
            "volume":  info.get('regularMarketVolume', 'N/A'),
            "source":  "yfinance"
        })
    except Exception as e:
        return jsonify({"error": str(e), "symbol": symbol}), 500


@app.route('/market/chart', methods=['GET'])
def market_chart():
    """OHLCV chart data for any NSE/NYSE stock via Twelve Data, falls back to yfinance."""
    symbol   = request.args.get('symbol', 'RELIANCE')
    exchange = request.args.get('exchange', 'NSE')
    interval = request.args.get('interval', '1day')   # 1min,5min,15min,1h,1day,1week
    outputsize = request.args.get('period', '90')      # number of data points
    try:
        if TWELVE_DATA_KEY:
            url = (f"https://api.twelvedata.com/time_series"
                   f"?symbol={symbol}&exchange={exchange}"
                   f"&interval={interval}&outputsize={outputsize}&apikey={TWELVE_DATA_KEY}")
            raw = requests.get(url, timeout=10).json()
            if raw.get('status') != 'error' and 'values' in raw:
                chart = []
                for item in reversed(raw['values']):
                    chart.append({
                        "date":   item['datetime'][:10],
                        "open":   float(item.get('open', 0)),
                        "high":   float(item.get('high', 0)),
                        "low":    float(item.get('low', 0)),
                        "price":  float(item.get('close', 0)),
                        "volume": int(float(item.get('volume', 0)))
                    })
                return jsonify(chart)
        # Fallback: yfinance
        yf_sym = f"{symbol}.NS" if exchange == 'NSE' else symbol
        hist = yf.Ticker(yf_sym).history(period="1y")
        chart = []
        for date, row in hist.iterrows():
            chart.append({
                "date":   date.strftime('%Y-%m-%d'),
                "open":   round(float(row['Open']), 2),
                "high":   round(float(row['High']), 2),
                "low":    round(float(row['Low']), 2),
                "price":  round(float(row['Close']), 2),
                "volume": int(row.get('Volume', 0))
            })
        return jsonify(chart)
    except Exception as e:
        return jsonify([]), 200  # Return empty array, not 500, so frontend handles gracefully


@app.route('/market/ipos', methods=['GET'])
def ipo_calendar():
    """Upcoming IPO calendar from Twelve Data + Finnhub combined, deduplicated."""
    from datetime import timedelta
    start = datetime.utcnow().strftime('%Y-%m-%d')
    end   = (datetime.utcnow() + timedelta(days=90)).strftime('%Y-%m-%d')
    results = []
    seen_names = set()

    # Source 1: Twelve Data
    if TWELVE_DATA_KEY:
        try:
            td_url = f"https://api.twelvedata.com/ipo_calendar?start_date={start}&end_date={end}&apikey={TWELVE_DATA_KEY}"
            td_data = requests.get(td_url, timeout=6).json()
            for ipo in td_data.get('ipos', []):
                name = ipo.get('name', '').strip()
                if name and name not in seen_names:
                    seen_names.add(name)
                    results.append({
                        "company":    name,
                        "exchange":   ipo.get('exchange', 'N/A'),
                        "date":       ipo.get('date', 'N/A'),
                        "price_low":  ipo.get('price_from', 'N/A'),
                        "price_high": ipo.get('price_to', 'N/A'),
                        "shares":     ipo.get('shares', 'N/A'),
                        "source":     "Twelve Data"
                    })
        except:
            pass

    # Source 2: Finnhub (catches more global listings)
    if FINNHUB_KEY:
        try:
            fh_url = f"https://finnhub.io/api/v1/calendar/ipo?from={start}&to={end}&token={FINNHUB_KEY}"
            fh_data = requests.get(fh_url, timeout=6).json()
            for ipo in fh_data.get('ipoCalendar', []):
                name = ipo.get('name', '').strip()
                if name and name not in seen_names:
                    seen_names.add(name)
                    results.append({
                        "company":    name,
                        "exchange":   ipo.get('exchange', 'N/A'),
                        "date":       ipo.get('date', 'N/A'),
                        "price_low":  ipo.get('price', 'N/A'),
                        "price_high": ipo.get('price', 'N/A'),
                        "shares":     ipo.get('numberOfShares', 'N/A'),
                        "source":     "Finnhub"
                    })
        except:
            pass

    results.sort(key=lambda x: x.get('date', ''))
    return jsonify(results)


@app.route('/market/movers', methods=['GET'])
def market_movers():
    """NSE top movers via yfinance fast_info (reliable) with batch fallback."""
    NIFTY20 = [
        'RELIANCE.NS','TCS.NS','HDFCBANK.NS','INFY.NS','ICICIBANK.NS',
        'SBIN.NS','BHARTIARTL.NS','WIPRO.NS','LT.NS','AXISBANK.NS',
        'MARUTI.NS','BAJFINANCE.NS','TECHM.NS','SUNPHARMA.NS','TITAN.NS',
        'NESTLEIND.NS','ULTRACEMCO.NS','ASIANPAINT.NS','KOTAKBANK.NS','HINDUNILVR.NS'
    ]
    results = []
    err_msg = None

    # ── Method 1: fast_info (single call per ticker, very fast) ────────
    try:
        for sym in NIFTY20:
            try:
                fi = yf.Ticker(sym).fast_info
                price = getattr(fi, 'last_price', None)
                prev  = getattr(fi, 'previous_close', None)
                if price and prev and prev > 0:
                    pct = round(((price - prev) / prev) * 100, 2)
                    results.append({
                        'symbol':  sym.replace('.NS', ''),
                        'ltp':     round(price, 2),
                        'pChange': pct
                    })
            except:
                continue
    except Exception as e:
        err_msg = f"fast_info failed: {e}"

    # ── Method 2 fallback: yf.download batch if fast_info yielded nothing ─
    if not results:
        try:
            raw = yf.download(NIFTY20, period='5d', interval='1d',
                              auto_adjust=True, progress=False)
            # Safely access Close DataFrame
            if hasattr(raw, 'columns') and len(raw.columns):
                # Handle both flat and MultiIndex columns
                if hasattr(raw.columns, 'levels'):
                    close_df = raw['Close']  # MultiIndex: (field, ticker)
                else:
                    close_df = raw            # Flat (single ticker edge case)

                for sym in NIFTY20:
                    try:
                        col = sym if sym in close_df.columns else None
                        if col is None:
                            continue
                        prices = close_df[col].dropna()
                        if len(prices) < 2:
                            continue
                        prev = float(prices.iloc[-2])
                        curr = float(prices.iloc[-1])
                        if prev == 0:
                            continue
                        pct = round(((curr - prev) / prev) * 100, 2)
                        results.append({
                            'symbol':  sym.replace('.NS', ''),
                            'ltp':     round(curr, 2),
                            'pChange': pct
                        })
                    except:
                        continue
        except Exception as e2:
            err_msg = (err_msg or '') + f" | batch failed: {e2}"

    results.sort(key=lambda x: x['pChange'], reverse=True)
    gainers = [r for r in results if r['pChange'] >= 0][:10]
    losers  = sorted([r for r in results if r['pChange'] < 0], key=lambda x: x['pChange'])[:10]
    return jsonify({
        'gainers': gainers,
        'losers':  losers,
        'indices': {},
        'error':   err_msg if not results else None
    })


@app.route('/market/index', methods=['GET'])
def market_index():
    """Index quote (Nifty50, S&P500, etc.) via Twelve Data, falls back to yfinance."""
    index    = request.args.get('index', '^NSEI')
    exchange = request.args.get('exchange', '')
    try:
        if TWELVE_DATA_KEY:
            url = f"https://api.twelvedata.com/quote?symbol={index}&apikey={TWELVE_DATA_KEY}"
            resp = requests.get(url, timeout=6).json()
            if resp.get('status') != 'error' and 'close' in resp:
                return jsonify({
                    "symbol":     resp.get('symbol', index),
                    "name":       resp.get('name', index),
                    "price":      resp.get('close', 'N/A'),
                    "change":     resp.get('change', 'N/A'),
                    "pct_change": resp.get('percent_change', 'N/A'),
                    "source":     "Twelve Data"
                })
    except Exception as e:
        return jsonify({"error": str(e), "symbol": index}), 500


# ── New Asset & Wealth Management, and Competitor Intelligence Endpoints ──

@app.route('/api/user/assets', methods=['GET'])
def get_user_assets():
    logged_user = get_logged_in_user()
    if not logged_user:
        return jsonify({"error": "Unauthorized"}), 401
    username = logged_user["username"]
    
    user_assets = list(assets_col.find({"username": username}))
    for asset in user_assets:
        asset["id"] = str(asset["_id"])
        del asset["_id"]
        ticker = asset.get("ticker", "").strip()
        if ticker:
            try:
                stock = yf.Ticker(ticker)
                price = None
                if hasattr(stock, 'fast_info') and getattr(stock.fast_info, 'last_price', None) is not None:
                    price = stock.fast_info.last_price
                else:
                    price = stock.info.get('currentPrice', None)
                if price is not None:
                    asset["current_price"] = round(price, 2)
                else:
                    asset["current_price"] = asset.get("purchase_price")
            except Exception as e:
                print(f"Error fetching live price for {ticker}: {e}")
                asset["current_price"] = asset.get("purchase_price")
        else:
            asset["current_price"] = asset.get("purchase_price")
    return jsonify(user_assets)

@app.route('/api/user/assets', methods=['POST'])
def add_user_asset():
    logged_user = get_logged_in_user()
    if not logged_user:
        return jsonify({"error": "Unauthorized"}), 401
    username = logged_user["username"]
    
    data = request.json
    name = data.get("name")
    ticker = data.get("ticker", "").strip()
    category = data.get("type", "Stock")
    quantity = float(data.get("quantity", 1))
    purchase_price = float(data.get("purchase_price", 0))
    purchase_date = data.get("purchase_date", datetime.utcnow().strftime("%Y-%m-%d"))
    
    asset_doc = {
        "username": username,
        "name": name,
        "ticker": ticker,
        "type": category,
        "quantity": quantity,
        "purchase_price": purchase_price,
        "purchase_date": purchase_date,
        "created_at": datetime.utcnow()
    }
    result = assets_col.insert_one(asset_doc)
    asset_doc["id"] = str(result.inserted_id)
    del asset_doc["_id"]
    return jsonify({"status": "success", "asset": asset_doc})

@app.route('/api/user/assets/delete', methods=['POST'])
def delete_user_asset():
    logged_user = get_logged_in_user()
    if not logged_user:
        return jsonify({"error": "Unauthorized"}), 401
    username = logged_user["username"]
    
    data = request.json
    asset_id = data.get("id")
    if not asset_id:
        return jsonify({"error": "Missing asset ID"}), 400
        
    from bson.objectid import ObjectId
    res = assets_col.delete_one({"_id": ObjectId(asset_id), "username": username})
    if res.deleted_count > 0:
        return jsonify({"status": "success"})
    return jsonify({"error": "Asset not found or unauthorized"}), 404

@app.route('/api/user/wealth/goals', methods=['GET'])
def get_user_goals():
    logged_user = get_logged_in_user()
    if not logged_user:
        return jsonify({"error": "Unauthorized"}), 401
    username = logged_user["username"]
    
    user_goals = list(wealth_goals_col.find({"username": username}))
    for goal in user_goals:
        goal["id"] = str(goal["_id"])
        del goal["_id"]
    return jsonify(user_goals)

@app.route('/api/user/wealth/goals', methods=['POST'])
def add_user_goal():
    logged_user = get_logged_in_user()
    if not logged_user:
        return jsonify({"error": "Unauthorized"}), 401
    username = logged_user["username"]
    
    data = request.json
    name = data.get("name")
    target_amount = float(data.get("target_amount", 0))
    target_year = int(data.get("target_year", datetime.utcnow().year + 10))
    current_savings = float(data.get("current_savings", 0))
    monthly_contribution = float(data.get("monthly_contribution", 0))
    risk_appetite = data.get("risk_appetite", "Moderate")
    
    goal_doc = {
        "username": username,
        "name": name,
        "target_amount": target_amount,
        "target_year": target_year,
        "current_savings": current_savings,
        "monthly_contribution": monthly_contribution,
        "risk_appetite": risk_appetite,
        "created_at": datetime.utcnow()
    }
    result = wealth_goals_col.insert_one(goal_doc)
    goal_doc["id"] = str(result.inserted_id)
    del goal_doc["_id"]
    return jsonify({"status": "success", "goal": goal_doc})

@app.route('/api/user/wealth/goals/delete', methods=['POST'])
def delete_user_goal():
    logged_user = get_logged_in_user()
    if not logged_user:
        return jsonify({"error": "Unauthorized"}), 401
    username = logged_user["username"]
    
    data = request.json
    goal_id = data.get("id")
    if not goal_id:
        return jsonify({"error": "Missing goal ID"}), 400
        
    from bson.objectid import ObjectId
    res = wealth_goals_col.delete_one({"_id": ObjectId(goal_id), "username": username})
    if res.deleted_count > 0:
        return jsonify({"status": "success"})
    return jsonify({"error": "Goal not found or unauthorized"}), 404

@app.route('/api/user/wealth/advisory', methods=['POST'])
def get_wealth_advisory():
    logged_user = get_logged_in_user()
    if not logged_user:
        return jsonify({"error": "Unauthorized"}), 401
    username = logged_user["username"]
    
    data = request.json
    risk_profile = data.get("risk_profile", "Moderate")
    
    assets = list(assets_col.find({"username": username}))
    goals = list(wealth_goals_col.find({"username": username}))
    
    assets_summary = ", ".join([f"{a.get('name')} ({a.get('type')}: {a.get('quantity')} units bought at {a.get('purchase_price')})" for a in assets])
    if not assets_summary:
        assets_summary = "No holdings currently registered."
        
    goals_summary = ", ".join([f"{g.get('name')} (Target: {g.get('target_amount')} by {g.get('target_year')}, savings: {g.get('current_savings')}, monthly SIP: {g.get('monthly_contribution')})" for g in goals])
    if not goals_summary:
        goals_summary = "No financial goals configured yet."
        
    prompt = f"""
    Act as a highly experienced wealth planner and alternative asset portfolio specialist.
    
    Analyze the client's current profile:
    - Risk Appetite: {risk_profile}
    - Current Asset Holdings: {assets_summary}
    - Financial Savings Goals: {goals_summary}
    
    Generate a detailed financial advisory report. It must be highly customized to their risk profile and goals.
    Provide standard strategic asset allocation based on risk:
    - Conservative: Equities 20%, Fixed Income 60%, Gold 10%, Cash 10%
    - Moderate: Equities 50%, Fixed Income 30%, Alternatives/Crypto 10%, Cash 10%
    - Aggressive: Equities 75%, Fixed Income 10%, Alternatives/Crypto 10%, Cash 5%
    
    Calculate or estimate compounding feasibility for their goals and give actionable guidance.
    
    Return ONLY a valid JSON object matching this structure:
    {{
      "recommendedAllocation": {{
        "Equities": 50,
        "Fixed Income": 30,
        "Alternatives": 10,
        "Cash": 10
      }},
      "goalFeasibility": "A detailed assessment of whether they will hit their goals based on target years, monthly SIPs, and current savings under historical average returns.",
      "strategicAdvice": [
        "Strategy point 1 (customized to their profile)",
        "Strategy point 2",
        "Strategy point 3",
        "Strategy point 4"
      ],
      "riskMitigation": [
        "Risk advice point 1",
        "Risk advice point 2"
      ]
    }}
    
    Return ONLY valid JSON. No markdown, no backticks, no explanations.
    """
    
    alloc_map = {
        "Conservative": {"Equities": 20, "Fixed Income": 60, "Alternatives": 10, "Cash": 10},
        "Moderate": {"Equities": 50, "Fixed Income": 30, "Alternatives": 10, "Cash": 10},
        "Aggressive": {"Equities": 75, "Fixed Income": 10, "Alternatives": 10, "Cash": 5}
    }
    fallback_alloc = alloc_map.get(risk_profile, alloc_map["Moderate"])
    fallback_advisory = {
        "recommendedAllocation": fallback_alloc,
        "goalFeasibility": "Based on a standard historical compounding rate of 8-12%, your wealth goals are partially on-track. Adding a consistent monthly SIP will accelerate compounding and improve probability of meeting your long-term target amounts.",
        "strategicAdvice": [
            f"Implement a structured rebalancing schedule to maintain your {risk_profile} allocation target.",
            "Diversify into global low-cost index ETFs to capture broader equity market growth.",
            "Establish a liquid emergency fund covering 6 months of living expenses before deploying capital to high-volatility alternative assets.",
            "Utilize tax-sheltered accounts (like PPF, 401k or ELSS depending on region) to optimize compounding efficiency."
        ],
        "riskMitigation": [
            "Maintain portfolio diversification; avoid over-allocating more than 5% of total wealth to speculative assets like cryptocurrency.",
            "Ensure asset maturities align with goal horizons to avoid forced liquidations during market downturns."
        ]
    }
    
    if not API_KEY or API_KEY == "sk-proj-your-real-api-key-goes-here...":
        return jsonify(fallback_advisory)
        
    is_openrouter = API_KEY.startswith("sk-or")
    is_gemini = API_KEY.startswith("AIzaSy")
    
    try:
        if is_gemini:
            url = f"https://generativelanguage.googleapis.com/v1/models/gemini-1.5-flash:generateContent?key={API_KEY}"
            payload = {
                "contents": [{
                    "parts": [{"text": prompt}]
                }]
            }
            response = requests.post(url, json=payload, timeout=10)
            res_json = response.json()
            
            if 'candidates' in res_json and res_json['candidates']:
                content = res_json['candidates'][0]['content']['parts'][0]['text']
                content = content.replace("```json", "").replace("```", "").strip()
                return jsonify(json.loads(content))
            return jsonify(fallback_advisory)
            
        url = "https://openrouter.ai/api/v1/chat/completions" if is_openrouter else "https://api.openai.com/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json"
        }
        if is_openrouter:
            headers["HTTP-Referer"] = "http://localhost:5174"
            headers["X-Title"] = "DB Advisory Platform"
            
        response = requests.post(
            url,
            headers=headers,
            json={
                "model": "openai/gpt-4o-mini" if is_openrouter else "gpt-4o-mini",
                "messages": [{"role": "user", "content": prompt}],
                "response_format": {"type": "json_object"}
            },
            timeout=10
        )
        content = response.json()['choices'][0]['message']['content']
        return jsonify(json.loads(content))
    except Exception as e:
        print(f"Advisory generation failed: {e}")
        return jsonify(fallback_advisory)

@app.route('/api/competitors', methods=['GET'])
def get_competitor_firms():
    logged_user = get_logged_in_user()
    if not logged_user:
        return jsonify({"error": "Unauthorized"}), 403
        
    firms = list(competitors_col.find().sort("globalRank", 1))
    for f in firms:
        f["id"] = str(f["_id"])
        del f["_id"]
    return jsonify(firms)

@app.route('/api/competitors/acquisitions', methods=['GET'])
def get_competitor_acquisitions():
    logged_user = get_logged_in_user()
    if not logged_user:
        return jsonify({"error": "Unauthorized"}), 403
        
    deals = list(acquisitions_col.find().sort("announcementDate", -1))
    for d in deals:
        d["id"] = str(d["_id"])
        del d["_id"]
        if d.get("buyerFirmId"):
            d["buyerFirmId"] = str(d["buyerFirmId"])
        if d.get("announcementDate"):
            d["announcementDate"] = d["announcementDate"].strftime("%Y-%m-%d")
            
    # Build macro statistics
    pipeline = [
        {"$group": {
            "_id": "$targetCategory",
            "totalValue": {"$sum": "$dealValue"},
            "count": {"$sum": 1}
        }}
    ]
    stats = list(acquisitions_col.aggregate(pipeline))
    stats_formatted = []
    for s in stats:
        stats_formatted.append({
            "category": s["_id"],
            "totalValue": s["totalValue"],
            "count": s["count"]
        })
        
    return jsonify({
        "acquisitions": deals,
        "macroTrends": stats_formatted
    })


if __name__ == "__main__":
    port = int(os.getenv("PORT", 5005))
    print(f"Backend starting on port {port}...")
    app.run(host="0.0.0.0", port=port)
